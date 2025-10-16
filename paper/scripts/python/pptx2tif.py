#!/usr/bin/env python3
# -*- coding: utf-8 -*-
# Timestamp: "2025-05-06 20:20:52 (ywatanabe)"
# File: /home/ywatanabe/proj/SciTex/manuscript/scripts/python/pptx2tif.py
# ----------------------------------------
import os
__FILE__ = (
    "./manuscript/scripts/python/pptx2tif.py"
)
__DIR__ = os.path.dirname(__FILE__)
# ----------------------------------------
"""
PowerPoint to TIF Conversion Utility

This script converts PowerPoint presentations (.pptx) to TIF images,
optimized for inclusion in scientific manuscripts.
"""

import argparse
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import List, Optional, Union

# Optional imports that might not be installed
LIBREOFFICE_AVAILABLE = False
PYTHON_PPT_AVAILABLE = False
PIL_AVAILABLE = False

try:
    # Check for python-pptx package
    from pptx import Presentation

    PYTHON_PPT_AVAILABLE = True
except ImportError:
    pass

try:
    # Check for PIL/Pillow
    from PIL import Image

    PIL_AVAILABLE = True
except ImportError:
    pass


def check_libreoffice_installed() -> bool:
    """Check if LibreOffice is installed."""
    try:
        result = subprocess.run(
            ["which", "libreoffice"],
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            text=True,
        )
        return result.returncode == 0
    except Exception:
        return False


def convert_pptx_to_tif_libreoffice(
    input_path: str,
    output_dir: Optional[str] = None,
    resolution: int = 300,
    verbose: bool = False,
) -> List[str]:
    """
    Convert a PowerPoint file to TIF using LibreOffice.

    Args:
        input_path: Path to the PowerPoint file
        output_dir: Directory to save output files (defaults to same directory as input)
        resolution: Image resolution in DPI
        verbose: Whether to print detailed information

    Returns:
        List of generated TIF file paths
    """
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"PowerPoint file not found: {input_path}")

    # Set output directory
    if output_dir is None:
        output_dir = os.path.dirname(os.path.abspath(input_path))
    elif not os.path.exists(output_dir):
        os.makedirs(output_dir)

    input_path = os.path.abspath(input_path)
    output_dir = os.path.abspath(output_dir)

    # Get the base name without extension
    base_name = os.path.splitext(os.path.basename(input_path))[0]

    # Create a temporary directory for conversion
    with tempfile.TemporaryDirectory() as temp_dir:
        # Convert to TIF using LibreOffice
        if verbose:
            print(f"Converting {input_path} to TIF using LibreOffice...")

        cmd = [
            "libreoffice",
            "--headless",
            "--convert-to",
            "tiff",
            f"--outdir",
            temp_dir,
            input_path,
        ]

        try:
            result = subprocess.run(
                cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True
            )

            if result.returncode != 0:
                raise RuntimeError(
                    f"LibreOffice conversion failed: {result.stderr}"
                )

            if verbose:
                print(result.stdout)

            # Find generated files
            tif_files = [
                os.path.join(temp_dir, f)
                for f in os.listdir(temp_dir)
                if f.lower().endswith((".tif", ".tiff"))
            ]

            if not tif_files:
                raise FileNotFoundError(
                    "No TIF files were generated during conversion"
                )

            # Move files to output directory with proper naming
            output_files = []
            for i, tif_file in enumerate(tif_files):
                # For single slide presentations, use the base name
                # For multi-slide presentations, append slide number
                if len(tif_files) == 1:
                    output_name = f"{base_name}.tif"
                else:
                    output_name = f"{base_name}_slide_{i+1}.tif"

                output_path = os.path.join(output_dir, output_name)

                # Copy file to output directory
                with open(tif_file, "rb") as contents_file, open(
                    output_path, "wb"
                ) as dst_file:
                    dst_file.write(contents_file.read())

                output_files.append(output_path)

                if verbose:
                    print(f"Saved: {output_path}")

            return output_files

        except Exception as e:
            raise RuntimeError(
                f"Error during LibreOffice conversion: {str(e)}"
            )


def convert_pptx_to_tif_python(
    input_path: str,
    output_dir: Optional[str] = None,
    resolution: int = 300,
    verbose: bool = False,
) -> List[str]:
    """
    Convert a PowerPoint file to TIF using python-pptx and PIL.

    Note: This method has limitations and may not work for all presentations.

    Args:
        input_path: Path to the PowerPoint file
        output_dir: Directory to save output files (defaults to same directory as input)
        resolution: Image resolution in DPI
        verbose: Whether to print detailed information

    Returns:
        List of generated TIF file paths
    """
    if not PYTHON_PPT_AVAILABLE:
        raise ImportError("python-pptx package is not installed")

    if not PIL_AVAILABLE:
        raise ImportError("PIL/Pillow package is not installed")

    if not os.path.exists(input_path):
        raise FileNotFoundError(f"PowerPoint file not found: {input_path}")

    # Set output directory
    if output_dir is None:
        output_dir = os.path.dirname(os.path.abspath(input_path))
    elif not os.path.exists(output_dir):
        os.makedirs(output_dir)

    # Get the base name without extension
    base_name = os.path.splitext(os.path.basename(input_path))[0]

    # Load the presentation
    if verbose:
        print(f"Opening PowerPoint file: {input_path}")

    prs = Presentation(input_path)
    output_files = []

    # Convert each slide
    for i, slide in enumerate(prs.slides):
        if verbose:
            print(f"Processing slide {i+1}/{len(prs.slides)}")

        # For multi-slide presentations, append slide number
        # For single slide presentations, use the base name
        if len(prs.slides) == 1:
            output_name = f"{base_name}.tif"
        else:
            output_name = f"{base_name}_slide_{i+1}.tif"

        output_path = os.path.join(output_dir, output_name)

        # This is a placeholder - python-pptx cannot directly render slides
        # We would need an additional library to render the slides

        # Instead, we'll just output a message
        print(
            f"Warning: Python-only conversion is limited. Slide {i+1} would be saved to {output_path}"
        )
        print(
            "For better results, please install LibreOffice or use the LibreOffice conversion method"
        )

        # Add to output_files even though we're not creating the file
        output_files.append(output_path)

    return output_files


def convert_pptx_to_tif(
    input_path: Union[str, Path],
    output_dir: Optional[Union[str, Path]] = None,
    method: str = "auto",
    resolution: int = 300,
    crop_whitespace: bool = True,
    margin: int = 30,
    verbose: bool = False,
) -> List[str]:
    """
    Convert a PowerPoint file to TIF using the best available method.

    Args:
        input_path: Path to the PowerPoint file
        output_dir: Directory to save output files (defaults to same directory as input)
        method: Conversion method ('libreoffice', 'python', or 'auto')
        resolution: Image resolution in DPI
        crop_whitespace: Whether to crop excess whitespace
        margin: Margin in pixels to add around the content area when cropping
        verbose: Whether to print detailed information

    Returns:
        List of generated TIF file paths
    """
    # Convert paths to strings if they're Path objects
    if isinstance(input_path, Path):
        input_path = str(input_path)

    if output_dir is not None and isinstance(output_dir, Path):
        output_dir = str(output_dir)

    # Determine the best method to use
    if method == "auto":
        if check_libreoffice_installed():
            method = "libreoffice"
            if verbose:
                print("Using LibreOffice for conversion")
        elif PYTHON_PPT_AVAILABLE and PIL_AVAILABLE:
            method = "python"
            if verbose:
                print(
                    "Using python-pptx for conversion (limited functionality)"
                )
        else:
            raise RuntimeError(
                "No suitable conversion method available. Please install LibreOffice or "
                "the python-pptx and Pillow packages."
            )

    # Perform the conversion
    if method == "libreoffice":
        output_files = convert_pptx_to_tif_libreoffice(
            input_path, output_dir, resolution, verbose
        )
    elif method == "python":
        output_files = convert_pptx_to_tif_python(
            input_path, output_dir, resolution, verbose
        )
    else:
        raise ValueError(f"Unknown conversion method: {method}")

    # Crop whitespace if requested
    if crop_whitespace and output_files:
        if verbose:
            print("\nCropping whitespace from generated images...")

        # Import crop_tif dynamically to avoid circular imports
        try:
            from crop_tif import crop_tif

            for tif_file in output_files:
                if verbose:
                    print(f"Cropping: {tif_file}")

                try:
                    crop_tif(tif_file, tif_file, margin, True, verbose)
                except Exception as e:
                    print(f"Warning: Failed to crop {tif_file}: {e}")

        except ImportError:
            print(
                "Warning: crop_tif module not available. Skipping whitespace cropping."
            )

    return output_files


def batch_convert_pptx_to_tif(
    directory: Union[str, Path],
    output_dir: Optional[Union[str, Path]] = None,
    method: str = "auto",
    resolution: int = 300,
    crop_whitespace: bool = True,
    margin: int = 30,
    recursive: bool = False,
    verbose: bool = False,
) -> List[str]:
    """
    Convert all PowerPoint files in a directory to TIF.

    Args:
        directory: Directory containing PowerPoint files
        output_dir: Directory to save output files (defaults to same as input)
        method: Conversion method ('libreoffice', 'python', or 'auto')
        resolution: Image resolution in DPI
        crop_whitespace: Whether to crop excess whitespace
        margin: Margin in pixels to add around the content area when cropping
        recursive: Whether to process subdirectories
        verbose: Whether to print detailed information

    Returns:
        List of generated TIF file paths
    """
    # Convert paths to strings if they're Path objects
    if isinstance(directory, Path):
        directory = str(directory)

    if output_dir is not None and isinstance(output_dir, Path):
        output_dir = str(output_dir)

    if not os.path.isdir(directory):
        raise ValueError(f"Directory not found: {directory}")

    # Get the list of PowerPoint files
    pptx_files = []
    if recursive:
        for root, _, filenames in os.walk(directory):
            for filename in filenames:
                if filename.lower().endswith((".ppt", ".pptx")):
                    pptx_files.append(os.path.join(root, filename))
    else:
        pptx_files = [
            os.path.join(directory, f)
            for f in os.listdir(directory)
            if f.lower().endswith((".ppt", ".pptx"))
        ]

    if not pptx_files:
        print(f"No PowerPoint files found in {directory}")
        return []

    # Process each file
    all_output_files = []
    for pptx_file in pptx_files:
        if verbose:
            print(f"\nProcessing: {pptx_file}")

        # Determine output directory
        if output_dir is None:
            file_output_dir = os.path.dirname(pptx_file)
        else:
            rel_path = os.path.relpath(os.path.dirname(pptx_file), directory)
            file_output_dir = os.path.join(output_dir, rel_path)

            # Create the directory if it doesn't exist
            if not os.path.exists(file_output_dir):
                os.makedirs(file_output_dir)

        # Convert the file
        try:
            output_files = convert_pptx_to_tif(
                pptx_file,
                file_output_dir,
                method,
                resolution,
                crop_whitespace,
                margin,
                verbose,
            )
            all_output_files.extend(output_files)
        except Exception as e:
            print(f"Error processing {pptx_file}: {e}")

    return all_output_files


def main():
    """Parse command-line arguments and execute the appropriate action."""
    # Set up argument parser
    parser = argparse.ArgumentParser(
        description="Convert PowerPoint files to TIF format.",
        formatter_class=argparse.ArgumentDefaultsHelpFormatter,
    )

    # Add subparsers for the different modes
    subparsers = parser.add_subparsers(dest="mode", help="Operation mode")

    # Single file mode
    file_parser = subparsers.add_parser("file", help="Process a single file")
    file_parser.add_argument(
        "-i", "--input", required=True, help="Input PowerPoint file path"
    )
    file_parser.add_argument(
        "-o", "--output-dir", help="Output directory for TIF files"
    )

    # Batch mode
    batch_parser = subparsers.add_parser(
        "batch", help="Process multiple files"
    )
    batch_parser.add_argument(
        "-d",
        "--directory",
        required=True,
        help="Directory containing PowerPoint files",
    )
    batch_parser.add_argument(
        "-o", "--output-dir", help="Output directory for TIF files"
    )
    batch_parser.add_argument(
        "-r",
        "--recursive",
        action="store_true",
        help="Process subdirectories recursively",
    )

    # Common arguments
    for subparser in [file_parser, batch_parser]:
        subparser.add_argument(
            "--method",
            choices=["auto", "libreoffice", "python"],
            default="auto",
            help="Conversion method to use",
        )
        subparser.add_argument(
            "--resolution",
            type=int,
            default=300,
            help="Output image resolution (DPI)",
        )
        subparser.add_argument(
            "--no-crop",
            action="store_true",
            help="Disable automatic cropping of whitespace",
        )
        subparser.add_argument(
            "--margin",
            type=int,
            default=30,
            help="Margin size around the content area when cropping",
        )
        subparser.add_argument(
            "-v",
            "--verbose",
            action="store_true",
            help="Enable verbose output",
        )

    # Parse arguments
    args = parser.parse_args()

    # Execute the appropriate action
    if args.mode == "file":
        try:
            output_files = convert_pptx_to_tif(
                args.input,
                args.output_dir,
                args.method,
                args.resolution,
                not args.no_crop,
                args.margin,
                args.verbose,
            )
            print(
                f"\nConversion complete. Generated {len(output_files)} TIF file(s)."
            )
        except Exception as e:
            print(f"Error: {e}")
            sys.exit(1)

    elif args.mode == "batch":
        try:
            output_files = batch_convert_pptx_to_tif(
                args.directory,
                args.output_dir,
                args.method,
                args.resolution,
                not args.no_crop,
                args.margin,
                args.recursive,
                args.verbose,
            )
            print(
                f"\nBatch conversion complete. Generated {len(output_files)} TIF file(s)."
            )
        except Exception as e:
            print(f"Error: {e}")
            sys.exit(1)

    else:
        parser.print_help()


if __name__ == "__main__":
    main()

# EOF
